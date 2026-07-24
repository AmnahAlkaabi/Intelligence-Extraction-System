"""Office Document Agent (L2, "Document" specialist alongside PDF) --
DOCX / PPTX via python-docx / python-pptx (pure-Python, no heavy deps).
"""
import asyncio
import logging

import docx
from pptx import Presentation

from app.models.schemas import FileCategory, ParsedDocument, TableBlock, TextBlock
from app.parsers.base import BaseParser

logger = logging.getLogger(__name__)


class OfficeParser(BaseParser):
    category = FileCategory.OFFICE

    async def parse(self, file_path: str) -> ParsedDocument:
        return await asyncio.to_thread(self._parse_sync, file_path)

    def _parse_sync(self, file_path: str) -> ParsedDocument:
        doc = ParsedDocument(source_file=file_path, category=self.category)
        try:
            if file_path.lower().endswith((".pptx",)):
                self._parse_pptx(file_path, doc)
            else:
                self._parse_docx(file_path, doc)
        except Exception as exc:  # noqa: BLE001
            logger.exception("Office parse failed on %s", file_path)
            doc.warnings.append(f"Office document parse error: {exc}")
        return doc

    def _parse_docx(self, file_path: str, doc: ParsedDocument) -> None:
        d = docx.Document(file_path)
        for para in d.paragraphs:
            if not para.text.strip():
                continue
            kind = "heading" if para.style and para.style.name and para.style.name.startswith("Heading") else "paragraph"
            doc.text_blocks.append(TextBlock(text=para.text, kind=kind))

        for table in d.tables:
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            headers = rows[0] if rows else []
            doc.tables.append(TableBlock(headers=headers, rows=rows[1:]))

        doc.metadata = {
            "paragraph_count": len(d.paragraphs),
            "table_count": len(d.tables),
            "parser": "python-docx",
        }
        if not doc.text_blocks and not doc.tables:
            doc.warnings.append("No extractable text or tables found in document.")

    def _parse_pptx(self, file_path: str, doc: ParsedDocument) -> None:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    doc.text_blocks.append(
                        TextBlock(text=shape.text_frame.text, page=i, kind="paragraph")
                    )
                if shape.has_table:
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    headers = rows[0] if rows else []
                    doc.tables.append(TableBlock(page=i, headers=headers, rows=rows[1:]))
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                doc.text_blocks.append(
                    TextBlock(text=slide.notes_slide.notes_text_frame.text, page=i, kind="notes")
                )

        doc.metadata = {"slide_count": len(prs.slides), "parser": "python-pptx"}
        if not doc.text_blocks and not doc.tables:
            doc.warnings.append("No extractable text or tables found in presentation.")
